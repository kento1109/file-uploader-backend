""" extract module from ppt """
import os
import shutil

import pandas as pd
import pptx
import spacy

nlp = spacy.load('ja_ginza')

STOP_WORDS = [
    'ため', 'の', '主', '年', '以来', 'ところ', 'よう', '方', 'どこ', 'こと', '者', 'さい', 'だれ', '感じ', '自体', '何', '.',
    '対応', '形', '人', '様々', '化'
]

# read excel file
DATAFRAMES = pd.read_excel('salary_table.xlsx', engine="openpyxl", sheet_name=None)
SALARY_TABLE = DATAFRAMES['salary']
KEYWORD_TABLE = DATAFRAMES['keyword']
KEYWORD_TABLE.fillna(0, inplace=True)


def tokenize(text):
    return " ".join([token.orth_ for token in nlp(text)])


def copy_file(file):
    """ copy file to current directory """
    with open(file.filename, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)


def delete_file(filename):
    """ delete file at current directory """
    os.remove(filename)


def get_presentation_obj(file):
    """ create presentation object """
    prs = pptx.Presentation(file)
    return prs


def read_text_box(slide):
    """ read text box """
    extracted_tokens = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            tokens = nlp_process(shape.text)
            extracted_tokens.extend(tokens)
    return extracted_tokens


def extract_key_tokens(text):
    doc = nlp(text)
    extracted_tokens = []
    for token in doc:
        if token.pos_ == 'NOUN':
            if token.lemma_ not in STOP_WORDS:
                extracted_tokens.append(token.lemma_)
    return extracted_tokens


def nlp_process(text):
    zen2han = str.maketrans('１２３４５６７８９０', '1234567890')
    text = text.translate(zen2han)
    text = text.replace('\n', '').replace('\x0b', '')
    tokens = extract_key_tokens(text)
    return tokens


def calculate_score(tokens):
    res = KEYWORD_TABLE[KEYWORD_TABLE['トークン'].isin(tokens)]
    occupation = res[['事務', '営業', 'エンジニア', '経営']].sum().idxmax() if len(res.index) else ''
    return SALARY_TABLE.iloc[0].to_dict().get(occupation, -1)
